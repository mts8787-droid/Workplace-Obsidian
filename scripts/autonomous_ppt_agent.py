import os
import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- Visual Fallback ---
try:
    import win32com.client
except ImportError:
    pass

def capture_slides(pptx_path: Path, output_dir: Path):
    print("📸 스크린 변환 캡쳐 진행 중...")
    if 'win32com' not in sys.modules:
        print("⚠️ win32com 모듈이 없어 캡쳐를 건너뜁니다.")
        return False
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        presentation = powerpoint.Presentations.Open(str(pptx_path), WithWindow=False)
        output_dir.mkdir(parents=True, exist_ok=True)
        presentation.Export(str(output_dir), "PNG")
        presentation.Close()
        powerpoint.Quit()
        return True
    except Exception as e:
        print(f"⚠️ 캡쳐 실패: {e}")
        return False

# --- Parsers ---
def parse_table(shape) -> str:
    if not hasattr(shape, "has_table") or not shape.has_table: return ""
    table = shape.table
    md_table = "\n"
    
    for row_idx, row in enumerate(table.rows):
        row_data = []
        for cell in row.cells:
            text = cell.text.replace("\n", " <br> ").replace("\r", "").strip()
            row_data.append(text if text else " ")
            
        md_table += "| " + " | ".join(row_data) + " |\n"
        if row_idx == 0:
            md_table += "|" + "|".join(["---"] * len(row.cells)) + "|\n"
            
    return md_table + "\n"

def parse_chart(shape) -> str:
    if not hasattr(shape, "has_chart") or not shape.has_chart: return ""
    chart = shape.chart
    chart_type = str(chart.chart_type).split('.')[-1]
    
    md_chart = f"\n> [!warning] 📊 **[차트 데이터 역산: {chart_type}]**\n"
    try:
        categories = [c.label for c in chart.plots[0].categories]
        md_chart += f"> | Series | " + " | ".join(categories) + " |\n"
        md_chart += f"> |---|{'---|'*len(categories)}\n"
        
        all_values = []
        for series in chart.series:
            values = [str(val) for val in series.values]
            md_chart += f"> | **{series.name}** | " + " | ".join(values) + " |\n"
            for val in series.values:
                if isinstance(val, (int, float)):
                    all_values.append((val, series.name))
                    
        if all_values:
            max_val = max(all_values, key=lambda x: x[0])
            md_chart += f">\n> 💡 **[분석]** 최고치 기록: **{max_val[1]} ({max_val[0]})**\n"
    except Exception as e:
        md_chart += f"> 데이터 추출 실패 ({e})\n"
        
    return md_chart + "\n"

def parse_diagram(shape) -> str:
    if shape.shape_type not in (6, 24): return "" 
    nodes = []
    if not hasattr(shape, "shapes"): return ""
        
    for sub_shape in shape.shapes:
        if sub_shape.has_text_frame:
            text = sub_shape.text.replace('\n', ' ').strip()
            if text:
                nodes.append({'text': text, 'top': sub_shape.top, 'left': sub_shape.left})
                
    if not nodes: return ""
    
    nodes.sort(key=lambda x: (x['top'] // 500000, x['left']))
    min_left = min(n['left'] for n in nodes)
    
    md_diagram = "\n> [!info] 🧩 **[도식 공간 분석]** 좌표 기반 흐름 (LLM Mermaid 렌더링 용)\n"
    for i, node in enumerate(nodes):
        depth_level = int((node['left'] - min_left) / 1000000) 
        indent = "  " * min(depth_level, 4)
        md_diagram += f"> {indent}{i+1}. {node['text']}\n"
        
    return md_diagram + "\n"

def parse_text(shape) -> str:
    if not hasattr(shape, "has_text_frame") or not shape.has_text_frame: return ""
    md_text = ""
    for paragraph in shape.text_frame.paragraphs:
        p_text = ""
        first_run_size = 12
        if paragraph.runs and paragraph.runs[0].font.size:
            first_run_size = paragraph.runs[0].font.size.pt
            
        for run in paragraph.runs:
            text = run.text.replace("\n", " ").strip()
            if not text: continue
            is_bold = run.font.bold
            color = str(run.font.color.rgb) if run.font.color and run.font.color.type == 1 else ""
            
            if color == "C00000": text = f"=={text}==" 
            elif color == "008000": text = f"*{text}*" 
            elif color == "FFF3F3": text = f"**[단점]** {text}" 
            if is_bold: text = f"**{text}**"
            p_text += text + " "
        
        p_text = p_text.strip()
        if not p_text: continue
            
        if "LGE Internal Use Only" in p_text or p_text.isdigit() or "/" in p_text and len(p_text) < 5:
            continue
            
        level = paragraph.level if hasattr(paragraph, 'level') else 0
        indent = "  " * level
        
        if first_run_size >= 20: md_text += f"## {p_text}\n\n"
        elif first_run_size >= 16: md_text += f"### 🎯 {p_text}\n\n"
        elif first_run_size >= 14: md_text += f"#### {p_text}\n\n"
        else: md_text += f"{indent}- {p_text}\n"
            
    return md_text + "\n" if md_text else ""

# --- Main Orchestrator ---
def process_ppt(pptx_path, out_dir_path):
    print("========================================")
    print(f"🚀 자율형 PPT 파싱 에이전트 시작")
    print(f"▶ 대상 파일: {pptx_path}")
    print("========================================")
    
    file_name = Path(pptx_path).stem
    output = [
        "---\n",
        f'title: "{file_name}"\n',
        'tags: ["PPT변환", "에이전트파싱", "자동화"]\n',
        "---\n\n",
        "👉 [원본 파일 열기](file:///" + pptx_path.replace("\\", "/") + ")\n\n"
    ]
    
    out_dir = Path(out_dir_path) / f"assets/{file_name}"
    capture_slides(Path(pptx_path), out_dir)
    
    prs = Presentation(pptx_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.replace('\n', ' ')
        else:
            title = f"Slide {slide_idx + 1}"
            
        output.append(f"# {slide_idx + 1}. {title}\n\n")
        output.append(f"![Slide {slide_idx + 1} Image](assets/{file_name}/슬라이드{slide_idx + 1}.PNG)\n\n")
        
        for shape in slide.shapes:
            if hasattr(shape, "has_table") and shape.has_table:
                output.append(parse_table(shape))
            elif hasattr(shape, "has_chart") and shape.has_chart:
                output.append(parse_chart(shape))
            else:
                try:
                    if getattr(shape, "shape_type", None) in (6, 24):
                        output.append(parse_diagram(shape))
                except NotImplementedError:
                    pass
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                output.append(parse_text(shape))
                
        output.append("\n\n")

    raw_markdown = "".join(output)
    
    final_markdown = raw_markdown

    md_out = str(Path(out_dir_path) / f"{file_name}.md")
    with open(md_out, "w", encoding="utf-8") as f:
        f.write(final_markdown)
    
    print("========================================")
    print(f"🎉 변환 완료: {md_out}")
    print("========================================")

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Autonomous PPT Parser Agent")
    parser.add_argument("ppt_path", type=str, help="Path to the PPTX file")
    parser.add_argument("--out", type=str, default="c:/Users/ts.moon/OneDrive - LG전자/obsidian/암묵지 변환 결과", help="Output directory")
    
    args = parser.parse_args()
    process_ppt(args.ppt_path, args.out)
